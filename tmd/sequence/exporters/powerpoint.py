"""
PowerPoint exporter for height map sequences.

This module provides functionality to export height map sequences to PowerPoint
presentations, with each frame as a slide.
"""

import os
import numpy as np
import logging
from typing import List, Optional, Union, Tuple

# Set up logger
logger = logging.getLogger(__name__)


def export_sequence_to_pptx(
    frames: List[np.ndarray],
    output_file: str,
    title: str = "Height Map Sequence",
    colormap: str = "terrain",
    include_frame_numbers: bool = True,
    dpi: int = 150
) -> Optional[str]:
    """
    Export a sequence of height maps to a PowerPoint presentation.
    
    Args:
        frames: List of 2D numpy arrays representing height maps
        output_file: Path to save the PowerPoint file
        title: Title for the presentation
        colormap: Matplotlib colormap name to use for rendering
        include_frame_numbers: Whether to include frame numbers on slides
        dpi: Resolution for the rendered images
        
    Returns:
        Path to the created file or None if failed
    """
    try:
        # Check for necessary packages
        try:
            from pptx import Presentation
            from pptx.util import Inches
            import matplotlib.pyplot as plt
            from matplotlib import cm
            import io
            from PIL import Image
        except ImportError as e:
            logger.error(f"Required package not found: {e}")
            logger.error("Please install python-pptx, matplotlib and Pillow packages")
            return None
            
        # Ensure frames list is not empty
        if not frames or len(frames) == 0:
            logger.error("No frames provided for PowerPoint export")
            return None
            
        # Ensure output directory exists
        os.makedirs(os.path.dirname(os.path.abspath(output_file)), exist_ok=True)
        
        # Create presentation
        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = f"{len(frames)} frames"
        
        # Add content slides
        content_slide_layout = prs.slide_layouts[5]  # Blank slide layout
        
        # Get colormap
        cmap = cm.get_cmap(colormap)
        
        # Process each frame
        for i, frame in enumerate(frames):
            # Create a new slide
            slide = prs.slides.add_slide(content_slide_layout)
            
            # Add frame number if requested
            if include_frame_numbers:
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(1), Inches(0.5))
                tf = txBox.text_frame
                tf.text = f"Frame {i+1}"
            
            # Render the height map as an image
            fig, ax = plt.subplots(figsize=(8, 6))
            
            # Normalize the frame data for consistent coloring
            norm_frame = (frame - np.min(frame)) / (np.max(frame) - np.min(frame) + 1e-10)
            
            # Display as an image
            ax.imshow(norm_frame, cmap=cmap)
            ax.axis('off')  # Hide axes
            
            # Save to memory buffer
            buf = io.BytesIO()
            fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight')
            buf.seek(0)
            
            # Close the matplotlib figure to free memory
            plt.close(fig)
            
            # Add the image to the slide
            img_path = io.BytesIO(buf.read())
            slide.shapes.add_picture(img_path, Inches(1), Inches(1), width=Inches(8))
        
        # Save the presentation
        prs.save(output_file)
        logger.info(f"PowerPoint presentation saved to {output_file}")
        
        return output_file
    
    except Exception as e:
        logger.error(f"Error exporting to PowerPoint: {e}")
        import traceback
        traceback.print_exc()
        return None
