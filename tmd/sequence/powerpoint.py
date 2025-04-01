class PowerPointExporter(BaseExporter):
    """
    Exporter for converting height map sequences into PowerPoint presentations.
    """
    
    def export(self, **kwargs) -> Optional[str]:
        """
        Expects the following kwargs:
            - frames: List of 2D numpy arrays (required)
            - output_file: Destination path for the PPTX file (required)
            - title: Presentation title (default: "Height Map Sequence")
            - colormap: Matplotlib colormap name (default: 'terrain')
            - include_frame_numbers: Whether to include frame numbers (default: True)
            - dpi: Resolution for rendered images (default: 150)
        """
        try:
            # Import dependencies
            from pptx import Presentation
            from pptx.util import Inches
            import matplotlib.pyplot as plt
            from matplotlib import cm
            import io
            from PIL import Image
            
            frames: List[np.ndarray] = kwargs.get('frames', [])
            output_file: str = kwargs.get('output_file')
            title: str = kwargs.get('title', "Height Map Sequence")
            colormap: str = kwargs.get('colormap', 'terrain')
            include_frame_numbers: bool = kwargs.get('include_frame_numbers', True)
            dpi: int = kwargs.get('dpi', 150)
            
            if not frames:
                logger.error("No frames provided for PowerPoint export")
                return None
            if not output_file:
                logger.error("Output file path must be provided for PPTX export")
                return None
            
            os.makedirs(os.path.dirname(os.path.abspath(output_file)), exist_ok=True)
            prs = Presentation()
            
            # Create title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = title
            slide.placeholders[1].text = f"{len(frames)} frames"
            
            content_slide_layout = prs.slide_layouts[5]
            cmap_obj = cm.get_cmap(colormap)
            
            for i, frame in enumerate(frames):
                slide = prs.slides.add_slide(content_slide_layout)
                
                if include_frame_numbers:
                    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(1), Inches(0.5))
                    txBox.text_frame.text = f"Frame {i+1}"
                
                fig, ax = plt.subplots(figsize=(8, 6))
                norm_frame = (frame - np.min(frame)) / (np.max(frame) - np.min(frame) + 1e-10)
                ax.imshow(norm_frame, cmap=cmap_obj)
                ax.axis('off')
                
                buf = io.BytesIO()
                fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight')
                buf.seek(0)
                plt.close(fig)
                
                slide.shapes.add_picture(buf, Inches(1), Inches(1), width=Inches(8))
            
            prs.save(output_file)
            logger.info(f"PowerPoint presentation saved to {output_file}")
            return output_file
        
        except Exception as e:
            logger.error(f"Error exporting to PowerPoint: {e}")
            return None